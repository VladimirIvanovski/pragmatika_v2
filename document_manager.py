#!/usr/bin/env python3
"""
Document Management System for the Intercultural Pragmatics Website
Handles scanning, processing, and metadata extraction from documents
"""

import os
import re
import json
from typing import Dict, List, Any, Optional
from datetime import datetime
import mimetypes
import hashlib

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False

try:
    import cloudinary
    import cloudinary.uploader
    CLOUDINARY_AVAILABLE = True
except ImportError:
    CLOUDINARY_AVAILABLE = False

DOCUMENTS_DIR = 'documents'
METADATA_FILE = 'content/documents_metadata.json'
EXTRACTED_IMAGES_DIR = 'static/extracted_images'
CLOUDINARY_CACHE_FILE = 'content/cloudinary_cache.json'

# ── Cloudinary setup ───────────────────────────────────────────────────────
_CLOUD_NAME = os.environ.get('CLOUDINARY_CLOUD_NAME')
_CLOUD_KEY   = os.environ.get('CLOUDINARY_API_KEY')
_CLOUD_SEC   = os.environ.get('CLOUDINARY_API_SECRET')
USE_CLOUDINARY = CLOUDINARY_AVAILABLE and all([_CLOUD_NAME, _CLOUD_KEY, _CLOUD_SEC])

if USE_CLOUDINARY:
    cloudinary.config(cloud_name=_CLOUD_NAME, api_key=_CLOUD_KEY, api_secret=_CLOUD_SEC)

def _load_cloudinary_cache() -> Dict[str, str]:
    """Load hash→URL cache from disk."""
    if os.path.exists(CLOUDINARY_CACHE_FILE):
        try:
            with open(CLOUDINARY_CACHE_FILE, 'r', encoding='utf-8') as f:
                return json.load(f)
        except Exception:
            pass
    return {}

def _save_cloudinary_cache(cache: Dict[str, str]) -> None:
    try:
        os.makedirs(os.path.dirname(CLOUDINARY_CACHE_FILE), exist_ok=True)
        with open(CLOUDINARY_CACHE_FILE, 'w', encoding='utf-8') as f:
            json.dump(cache, f, indent=2)
    except OSError:
        pass

def _store_image(image_bytes: bytes, public_id_hint: str) -> str:
    """Store image bytes. Returns a URL (Cloudinary) or local path.
    
    Uses the MD5 hash as the Cloudinary public_id for deduplication.
    Falls back to local disk when Cloudinary is not configured.
    """
    file_hash = hashlib.md5(image_bytes).hexdigest()

    if USE_CLOUDINARY:
        cache = _load_cloudinary_cache()
        if file_hash in cache:
            return cache[file_hash]
        try:
            result = cloudinary.uploader.upload(
                image_bytes,
                public_id=f'pragmatika/{file_hash}',
                overwrite=False,
                resource_type='image',
                format='webp',
            )
            url = result['secure_url']
            cache[file_hash] = url
            _save_cloudinary_cache(cache)
            return url
        except Exception as e:
            print(f'Cloudinary upload failed: {e}')

    # Local fallback (skip on read-only filesystems)
    short_hash = file_hash[:8]
    clean_hint = "".join(c for c in public_id_hint if c.isalnum() or c in (' ', '-', '_')).strip()
    image_filename = f'{clean_hint}_{short_hash}.png'
    image_path = os.path.join(EXTRACTED_IMAGES_DIR, image_filename)
    try:
        os.makedirs(EXTRACTED_IMAGES_DIR, exist_ok=True)
        with open(image_path, 'wb') as f:
            f.write(image_bytes)
    except OSError:
        return ''
    return f'/static/extracted_images/{image_filename}'

def ensure_documents_dir():
    """Ensure documents directory exists (no-op on read-only filesystems like Vercel)."""
    try:
        os.makedirs(DOCUMENTS_DIR, exist_ok=True)
        os.makedirs(os.path.dirname(METADATA_FILE), exist_ok=True)
        os.makedirs(EXTRACTED_IMAGES_DIR, exist_ok=True)
    except OSError:
        pass

def get_file_metadata(file_path: str) -> Dict[str, Any]:
    """Extract metadata from a file"""
    file_stat = os.stat(file_path)
    file_name = os.path.basename(file_path)
    file_ext = os.path.splitext(file_name)[1].lower()
    
    metadata = {
        'filename': file_name,
        'path': file_path,
        'size': file_stat.st_size,
        'size_mb': round(file_stat.st_size / (1024 * 1024), 2),
        'modified': datetime.fromtimestamp(file_stat.st_mtime).isoformat(),
        'type': file_ext[1:] if file_ext else 'unknown',
        'mime_type': mimetypes.guess_type(file_path)[0] or 'application/octet-stream'
    }
    
    # Extract additional metadata based on file type
    if file_ext == '.docx' and DOCX_AVAILABLE:
        try:
            doc = DocxDocument(file_path)
            metadata['title'] = file_name.replace('.docx', '')
            metadata['page_count'] = len(doc.paragraphs)
            metadata['word_count'] = sum(len(p.text.split()) for p in doc.paragraphs)
            metadata['description'] = doc.paragraphs[0].text[:200] if doc.paragraphs else ''
            metadata['has_images'] = len(doc.inline_shapes) > 0
        except Exception as e:
            metadata['error'] = str(e)
    
    elif file_ext == '.pptx' and PPTX_AVAILABLE:
        try:
            prs = Presentation(file_path)
            metadata['title'] = file_name.replace('.pptx', '')
            metadata['slide_count'] = len(prs.slides)
            metadata['description'] = 'PowerPoint презентација'
            metadata['has_images'] = any(len(slide.shapes) > 0 for slide in prs.slides)
        except Exception as e:
            metadata['error'] = str(e)
    
    elif file_ext == '.pdf':
        metadata['title'] = file_name.replace('.pdf', '')
        metadata['description'] = 'PDF документ'
    
    return metadata

def scan_documents() -> List[Dict[str, Any]]:
    """Scan documents directory recursively and return list of all documents with metadata"""
    ensure_documents_dir()
    documents = []
    
    if not os.path.exists(DOCUMENTS_DIR):
        return documents
    
    supported_extensions = ['.pdf', '.docx', '.pptx', '.doc', '.ppt']
    
    for root, _dirs, files in os.walk(DOCUMENTS_DIR):
        rel_root = os.path.relpath(root, DOCUMENTS_DIR)
        folder = '' if rel_root == '.' else rel_root.split(os.sep)[0]
        
        for file_name in files:
            file_ext = os.path.splitext(file_name)[1].lower()
            if file_ext in supported_extensions:
                file_path = os.path.join(root, file_name)
                try:
                    metadata = get_file_metadata(file_path)
                    metadata['folder'] = folder
                    if rel_root == '.':
                        metadata['subfolder_path'] = file_name
                    else:
                        metadata['subfolder_path'] = os.path.join(rel_root, file_name).replace(os.sep, '/')
                    documents.append(metadata)
                except Exception as e:
                    print(f"Error processing {file_name}: {e}")
                    continue
    
    documents.sort(key=lambda x: x.get('modified', ''), reverse=True)
    return documents


def scan_documents_by_folder() -> Dict[str, List[Dict[str, Any]]]:
    """Return documents grouped by subfolder name.
    
    Returns a dict like {'Moduli': [...], 'Conferences': [...], '': [...top-level...]}
    """
    grouped: Dict[str, List[Dict[str, Any]]] = {}
    for doc in scan_documents():
        folder = doc.get('folder', '')
        grouped.setdefault(folder, []).append(doc)
    return grouped

IMAGE_EXTENSIONS = {'.jpg', '.jpeg', '.png', '.gif', '.webp'}

def _natural_sort_key(s: str):
    """Sort key for 'pic 1', 'pic 2', ... 'pic 10' to order 1,2,...,10."""
    parts = re.split(r'(\d+)', s.lower())
    return [int(p) if p.isdigit() else p for p in parts]

def get_conference_subfolder_images(subfolder_name: str) -> List[Dict[str, str]]:
    """Get image paths from a conference subfolder. Uploads to Cloudinary when available."""
    from urllib.parse import quote
    folder_path = os.path.join(DOCUMENTS_DIR, 'Conferences', subfolder_name)
    if not os.path.isdir(folder_path):
        return []
    images = []
    files = [f for f in os.listdir(folder_path) if os.path.splitext(f)[1].lower() in IMAGE_EXTENSIONS]
    for f in sorted(files, key=_natural_sort_key):
        img_path = os.path.join(folder_path, f)
        try:
            with open(img_path, 'rb') as fh:
                img_bytes = fh.read()
            hint = os.path.splitext(f)[0]
            url = _store_image(img_bytes, f'{subfolder_name}_{hint}')
            if url:
                images.append({'path': url})
                continue
        except Exception:
            pass
        rel = subfolder_name + '/' + f
        images.append({'path': '/documents/conference-image/' + quote(rel)})
    return images

def _extract_entry_title(content: dict, fallback: str) -> str:
    """Pick the best title from doc content: first Heading paragraph, or first non-empty paragraph."""
    for p in content.get('paragraphs') or []:
        style = p.get('style') or ''
        text = p.get('text', '').strip()
        if text and ('Heading' in style or 'Title' in style):
            return text[:200]
    for p in content.get('paragraphs') or []:
        text = p.get('text', '').strip()
        if text:
            return text[:200]
    return fallback

def get_conference_entries() -> List[Dict[str, Any]]:
    """Build conference entries: standalone docs + subfolders (doc + folder images)."""
    entries = []
    seen_subfolders = set()
    for doc in scan_documents():
        if doc.get('folder') != 'Conferences':
            continue
        subfolder_path = doc.get('subfolder_path', '')
        parts = subfolder_path.replace('/', os.sep).split(os.sep)
        is_subfolder = len(parts) >= 2
        subfolder = parts[1] if is_subfolder else None
        if is_subfolder and subfolder and subfolder not in seen_subfolders:
            seen_subfolders.add(subfolder)
            content = get_document_content(subfolder_path)
            if not isinstance(content, dict) or 'error' in content:
                content = {'paragraphs': [], 'images': [], 'tables': []}
            folder_images = get_conference_subfolder_images(subfolder)
            if folder_images:
                content = dict(content)
                content['images'] = folder_images + (content.get('images') or [])
            title = _extract_entry_title(content, subfolder)
            paras = content.get('paragraphs') or []
            if paras and paras[0].get('text', '').strip() == title.strip():
                content = dict(content)
                content['paragraphs'] = paras[1:]
            entries.append({'filename': doc['filename'], 'title': title, 'content': content})
        elif not is_subfolder:
            content = get_document_content(subfolder_path)
            title = _extract_entry_title(content, os.path.splitext(doc['filename'])[0]) if isinstance(content, dict) else os.path.splitext(doc['filename'])[0]
            entries.append({'filename': doc['filename'], 'title': title, 'content': content})
    return entries

def get_documents_by_type(doc_type: str = None) -> List[Dict[str, Any]]:
    """Get documents filtered by type"""
    all_docs = scan_documents()
    if doc_type:
        return [doc for doc in all_docs if doc.get('type', '').lower() == doc_type.lower()]
    return all_docs

def group_documents_by_type(documents: List[Dict[str, Any]]) -> Dict[str, List[Dict[str, Any]]]:
    """Group documents by their type"""
    grouped = {}
    for doc in documents:
        doc_type = doc.get('type', 'other')
        if doc_type not in grouped:
            grouped[doc_type] = []
        grouped[doc_type].append(doc)
    return grouped

def save_metadata(documents: List[Dict[str, Any]]) -> None:
    """Save document metadata to JSON file"""
    ensure_documents_dir()
    metadata = {
        'last_scan': datetime.now().isoformat(),
        'total_documents': len(documents),
        'documents': documents
    }
    with open(METADATA_FILE, 'w', encoding='utf-8') as f:
        json.dump(metadata, f, ensure_ascii=False, indent=2)

def load_metadata() -> Dict[str, Any]:
    """Load document metadata from JSON file"""
    if os.path.exists(METADATA_FILE):
        with open(METADATA_FILE, 'r', encoding='utf-8') as f:
            return json.load(f)
    return {'last_scan': None, 'total_documents': 0, 'documents': []}

def get_document_path(filename: str) -> Optional[str]:
    """Get full path to a document by filename.
    
    Accepts a bare filename (backward-compatible) or a subfolder-relative path
    like 'Moduli/file.pptx'. Falls back to a recursive search when a direct
    match isn't found.
    """
    direct = os.path.join(DOCUMENTS_DIR, filename)
    if os.path.isfile(direct):
        return direct
    
    target = os.path.basename(filename)
    for root, _dirs, files in os.walk(DOCUMENTS_DIR):
        if target in files:
            return os.path.join(root, target)
    return None

def extract_docx_content(file_path: str) -> Dict[str, Any]:
    """Extract text and images from a Word document"""
    if not DOCX_AVAILABLE:
        return {'error': 'python-docx not available'}
    
    try:
        doc = DocxDocument(file_path)
        content = {
            'title': os.path.splitext(os.path.basename(file_path))[0],
            'paragraphs': [],
            'images': []
        }
        
        # Extract paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                content['paragraphs'].append({
                    'text': para.text,
                    'style': para.style.name if para.style else None
                })
        
        # Extract tables
        content['tables'] = []
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text for cell in row.cells]
                table_data.append(row_data)
            if table_data:
                content['tables'].append(table_data)
        
        # Extract images
        image_index = 0
        filename_base = os.path.splitext(os.path.basename(file_path))[0]
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                try:
                    image_bytes = rel.target_part.blob
                    hint = f'{filename_base}_{image_index}'
                    url = _store_image(image_bytes, hint)
                    content['images'].append({'path': url, 'index': image_index})
                    image_index += 1
                except Exception as e:
                    print(f"Error extracting image: {e}")
                    continue
        
        return content
    except Exception as e:
        return {'error': str(e)}

def extract_pptx_content(file_path: str) -> Dict[str, Any]:
    """Extract slides and images from a PowerPoint presentation"""
    if not PPTX_AVAILABLE:
        return {'error': 'python-pptx not available'}
    
    try:
        prs = Presentation(file_path)
        content = {
            'title': os.path.splitext(os.path.basename(file_path))[0],
            'slides': []
        }
        
        slide_index = 0
        for slide in prs.slides:
            slide_content = {
                'number': slide_index + 1,
                'texts': [],
                'images': []
            }
            
            # Extract text and images from shapes
            for shape in slide.shapes:
                # Extract text
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content['texts'].append(shape.text.strip())
                
                # Extract images - try multiple methods
                image_extracted = False
                
                # Method 1: Direct image attribute
                if hasattr(shape, "image"):
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        image_extracted = True
                    except:
                        pass
                
                # Method 2: Check if shape has part with image
                if not image_extracted and hasattr(shape, "part"):
                    try:
                        for rel in shape.part.rels.values():
                            if "image" in rel.target_ref:
                                image_part = rel.target_part
                                image_bytes = image_part.blob
                                image_extracted = True
                                break
                    except:
                        pass
                
                # Store image if found
                if image_extracted:
                    try:
                        filename_base = os.path.splitext(os.path.basename(file_path))[0]
                        hint = f'{filename_base}_slide{slide_index}_img{len(slide_content["images"])}'
                        url = _store_image(image_bytes, hint)
                        slide_content['images'].append({
                            'path': url,
                            'alt': f'Slide {slide_index + 1} image {len(slide_content["images"]) + 1}'
                        })
                    except Exception as e:
                        print(f"Error storing image from slide {slide_index}: {e}")
                        continue
            
            if slide_content['texts'] or slide_content['images']:
                content['slides'].append(slide_content)
            
            slide_index += 1
        
        return content
    except Exception as e:
        return {'error': str(e)}

def get_document_content(filename: str) -> Dict[str, Any]:
    """Get extracted content from a document"""
    file_path = get_document_path(filename)
    if not file_path:
        return {'error': 'File not found'}
    
    file_ext = os.path.splitext(filename)[1].lower()
    
    if file_ext == '.docx':
        return extract_docx_content(file_path)
    elif file_ext == '.pptx':
        return extract_pptx_content(file_path)
    elif file_ext == '.pdf':
        # For PDF, return metadata and embed info
        return {
            'title': os.path.splitext(filename)[0],
            'type': 'pdf',
            'file_path': f'/documents/view/{filename}'
        }
    else:
        return {'error': 'Unsupported file type'}

if __name__ == '__main__':
    # Test the document manager
    import sys
    import io
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
    
    print("Scanning documents...")
    docs = scan_documents()
    print(f"Found {len(docs)} documents")
    for doc in docs:
        print(f"\n{doc['filename']}")
        print(f"  Type: {doc['type']}")
        print(f"  Size: {doc['size_mb']} MB")
        if 'slide_count' in doc:
            print(f"  Slides: {doc['slide_count']}")
        if 'word_count' in doc:
            print(f"  Words: {doc['word_count']}")
    
    save_metadata(docs)
    print(f"\nMetadata saved to {METADATA_FILE}")

